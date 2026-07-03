from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = r"C:\Users\금정산2-PC02\p2-spring\spring-ai-study"
OUT_MAIN = ROOT + r"\DAY1_DAY2_VISUAL.pptx"
OUT_STYLE = ROOT + r"\DAY1_DAY2_VISUAL_PALANTIR_STYLE.pptx"

W = Inches(13.333)
H = Inches(7.5)

BLACK = RGBColor(5, 8, 13)
PANEL = RGBColor(11, 18, 27)
PANEL_2 = RGBColor(16, 27, 41)
GRID = RGBColor(36, 52, 70)
TEXT = RGBColor(236, 242, 247)
MUTED = RGBColor(145, 164, 184)
CYAN = RGBColor(0, 221, 204)
BLUE = RGBColor(70, 136, 255)
AMBER = RGBColor(255, 192, 87)
RED = RGBColor(255, 93, 93)
GREEN = RGBColor(89, 232, 157)

FONT = "Malgun Gothic"


def set_text(run, size=16, color=TEXT, bold=False):
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold


def bg(slide):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BLACK
    for i in range(16):
        x = Inches(0.4 + i * 0.82)
        line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x, Inches(0.35), x, Inches(7.15))
        line.line.color.rgb = GRID
        line.line.transparency = 55
        line.line.width = Pt(0.35)
    for i in range(8):
        y = Inches(0.55 + i * 0.82)
        line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(0.25), y, Inches(13.05), y)
        line.line.color.rgb = GRID
        line.line.transparency = 55
        line.line.width = Pt(0.35)
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), W, Inches(0.06))
    accent.fill.solid()
    accent.fill.fore_color.rgb = CYAN
    accent.line.fill.background()


def title(slide, main, sub=None, tag="SPRING AI"):
    box = slide.shapes.add_textbox(Inches(0.55), Inches(0.42), Inches(9.6), Inches(0.65))
    p = box.text_frame.paragraphs[0]
    p.text = main
    set_text(p.runs[0], 25, TEXT, True)
    if sub:
        sbox = slide.shapes.add_textbox(Inches(0.58), Inches(1.1), Inches(9.8), Inches(0.38))
        sp = sbox.text_frame.paragraphs[0]
        sp.text = sub
        set_text(sp.runs[0], 11, MUTED, False)
    tag_box = slide.shapes.add_textbox(Inches(10.45), Inches(0.48), Inches(2.25), Inches(0.3))
    tp = tag_box.text_frame.paragraphs[0]
    tp.text = tag
    tp.alignment = PP_ALIGN.RIGHT
    set_text(tp.runs[0], 10, CYAN, True)


def panel(slide, x, y, w, h, label=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = PANEL
    shape.line.color.rgb = GRID
    shape.line.width = Pt(0.7)
    if label:
        label_box = slide.shapes.add_textbox(Inches(x + 0.18), Inches(y + 0.12), Inches(w - 0.3), Inches(0.28))
        lp = label_box.text_frame.paragraphs[0]
        lp.text = label.upper()
        set_text(lp.runs[0], 8, CYAN, True)
    return shape


def small_label(slide, x, y, text, color=MUTED):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(2.4), Inches(0.22))
    p = box.text_frame.paragraphs[0]
    p.text = text
    set_text(p.runs[0], 8, color, True)


def node(slide, x, y, text, color=CYAN, w=1.55, h=0.62):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = PANEL_2
    s.line.color.rgb = color
    s.line.width = Pt(1.2)
    p = s.text_frame.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.CENTER
    set_text(p.runs[0], 10, TEXT, True)
    return s


def connect(slide, a, b, color=GRID):
    x1 = a.left + a.width
    y1 = a.top + a.height // 2
    x2 = b.left
    y2 = b.top + b.height // 2
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    line.line.color.rgb = color
    line.line.width = Pt(1.2)


def metric(slide, x, y, k, v, color=CYAN):
    panel(slide, x, y, 2.25, 0.95)
    kp = slide.shapes.add_textbox(Inches(x + 0.17), Inches(y + 0.13), Inches(1.9), Inches(0.22))
    p = kp.text_frame.paragraphs[0]
    p.text = k
    set_text(p.runs[0], 8, MUTED, True)
    vp = slide.shapes.add_textbox(Inches(x + 0.17), Inches(y + 0.43), Inches(1.9), Inches(0.36))
    p2 = vp.text_frame.paragraphs[0]
    p2.text = v
    set_text(p2.runs[0], 18, color, True)


def bullet_panel(slide, x, y, w, h, header, bullets):
    panel(slide, x, y, w, h, header)
    body = slide.shapes.add_textbox(Inches(x + 0.28), Inches(y + 0.58), Inches(w - 0.5), Inches(h - 0.75))
    tf = body.text_frame
    tf.clear()
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = b
        p.space_after = Pt(7)
        set_text(p.runs[0], 15, TEXT, False)


def make_deck():
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H

    # 1
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide)
    title(slide, "Spring AI 학습 작전도", "Day1 첫 호출에서 Day2 구조화 응답까지", "DAY1-DAY2")
    metric(slide, 0.75, 1.75, "SCOPE", "2 DAYS", CYAN)
    metric(slide, 3.2, 1.75, "CORE FLOW", "C-S-C", BLUE)
    metric(slide, 5.65, 1.75, "OUTPUTS", "6 APIs", GREEN)
    metric(slide, 8.1, 1.75, "NEXT", "PROMPT OPS", AMBER)
    bullet_panel(slide, 0.75, 3.0, 5.8, 2.85, "핵심 메시지", [
        "Day1은 AI를 Spring 앱에 연결한 날",
        "Day2는 답변 형식을 통제하기 시작한 날",
        "이제 ch03-prompt에서 프롬프트 전략을 비교한다",
    ])
    bullet_panel(slide, 6.8, 3.0, 5.75, 2.85, "현재 위치", [
        "Controller, Service, ChatClient 흐름 이해",
        "String, Object, List 응답 처리 경험",
        "URL, Gradle, Git까지 실전 문제 정리",
    ])

    # 2
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide)
    title(slide, "학습 로드맵", "기초 복습에서 프롬프트 기법 분석까지", "ROADMAP")
    n1 = node(slide, 0.75, 3.0, "기초 복습", CYAN)
    n2 = node(slide, 2.55, 3.0, "Day1\n첫 호출", BLUE)
    n3 = node(slide, 4.35, 3.0, "Controller\nService", GREEN)
    n4 = node(slide, 6.15, 3.0, "Day2\nPrompt", CYAN)
    n5 = node(slide, 7.95, 3.0, "Structured\nOutput", AMBER)
    n6 = node(slide, 9.75, 3.0, "환경/URL\n디버깅", RED)
    n7 = node(slide, 11.55, 3.0, "ch03\nPrompt", CYAN)
    for a, b in [(n1, n2), (n2, n3), (n3, n4), (n4, n5), (n5, n6), (n6, n7)]:
        connect(slide, a, b)
    bullet_panel(slide, 1.0, 4.55, 11.3, 1.35, "읽는 법", [
        "왼쪽은 배운 순서, 오른쪽은 앞으로 확장할 주제입니다.",
        "중앙의 Controller -> Service -> ChatClient가 모든 예제의 공통 뼈대입니다.",
    ])

    # 3
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide)
    title(slide, "공통 실행 흐름", "모든 API가 같은 통로를 지나갑니다", "FLOW")
    a = node(slide, 0.9, 3.05, "User\nBrowser", CYAN, 1.45, 0.78)
    b = node(slide, 2.75, 3.05, "Controller\nRequest", BLUE, 1.55, 0.78)
    c = node(slide, 4.72, 3.05, "Service\nPrompt", GREEN, 1.55, 0.78)
    d = node(slide, 6.72, 3.05, "ChatClient\nCall", CYAN, 1.55, 0.78)
    e = node(slide, 8.72, 3.05, "AI Model\nGenerate", AMBER, 1.55, 0.78)
    f = node(slide, 10.72, 3.05, "Response\nReturn", GREEN, 1.55, 0.78)
    for x, y in [(a, b), (b, c), (c, d), (d, e), (e, f)]:
        connect(slide, x, y)
    small_label(slide, 1.0, 4.15, "0ms 요청", CYAN)
    small_label(slide, 2.82, 4.15, "10ms 파라미터", BLUE)
    small_label(slide, 4.78, 4.15, "20ms 프롬프트", GREEN)
    small_label(slide, 6.86, 4.15, "60ms 호출", CYAN)
    small_label(slide, 8.82, 4.15, "200ms+ 생성", AMBER)
    small_label(slide, 10.85, 4.15, "결과 반환", GREEN)

    # 4
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide)
    title(slide, "Day1", "AI를 앱에 연결하는 최소 단위", "FIRST CALL")
    bullet_panel(slide, 0.75, 1.75, 5.8, 4.75, "배운 것", [
        "@RestController: 요청을 받는 창구",
        "@Service: 실제 작업 담당",
        "@RequestParam: URL 값 꺼내기",
        "prompt() -> user() -> call() -> content()",
    ])
    bullet_panel(slide, 6.8, 1.75, 5.75, 4.75, "Python으로 치면", [
        "Flask route가 Controller 역할",
        "일반 함수가 Service 역할",
        "request.args.get(...)이 @RequestParam 역할",
        "AI 호출 함수가 ChatClient 역할",
    ])

    # 5
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide)
    title(slide, "Day2", "질문을 설계하고 응답 모양을 통제하기", "PROMPT OUTPUT")
    bullet_panel(slide, 0.75, 1.75, 5.8, 4.75, "Prompt 설계", [
        "그냥 질문하기에서 규칙 주기 단계로 이동",
        "{text}, {audience}, {mood} 같은 자리표시자 사용",
        "param()으로 실제 값을 넣음",
        "AI에게 대상, 길이, 형식을 함께 지시",
    ])
    bullet_panel(slide, 6.8, 1.75, 5.75, 4.75, "Structured Output", [
        "content(): 문자열 그대로 받기",
        "entity(): 객체로 변환해 받기",
        "record: 응답 모양 정의",
        "List는 ParameterizedTypeReference로 타입 전달",
    ])

    # 6
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide)
    title(slide, "기능 매트릭스", "Day2에서 만든 API와 응답 형태", "MATRIX")
    rows = [
        ("요약", "/api/summary", "String", CYAN),
        ("분류", "/api/classify", "String", BLUE),
        ("분류 JSON", "/api/classify/json", "JSON text", AMBER),
        ("분류 객체", "/api/classify/object", "InquiryResult", GREEN),
        ("영화 추천", "/api/movie", "List<MovieResponse>", GREEN),
        ("준비물 추천", "/api/packing", "List<String>", CYAN),
    ]
    y = 1.7
    for name, path, out, color in rows:
        panel(slide, 0.8, y, 11.7, 0.62)
        for x, text, w, c in [(1.05, name, 2.0, color), (3.1, path, 3.0, TEXT), (6.35, out, 2.9, TEXT)]:
            box = slide.shapes.add_textbox(Inches(x), Inches(y + 0.14), Inches(w), Inches(0.26))
            p = box.text_frame.paragraphs[0]
            p.text = text
            set_text(p.runs[0], 12, c, True if x == 1.05 else False)
        y += 0.75

    # 7
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide)
    title(slide, "디버깅에서 배운 것", "코드보다 실행 상태가 문제인 경우도 많았습니다", "INCIDENTS")
    bullet_panel(slide, 0.75, 1.65, 3.75, 4.9, "실행", [
        "day01이 8080에 떠 있는데 day02 API 호출",
        "IntelliJ 실행 대상 확인 필요",
        "JDK 21과 Gradle JVM 정렬",
    ])
    bullet_panel(slide, 4.8, 1.65, 3.75, 4.9, "URL", [
        "/api/classify/object",
        "/api/classifyobject는 다른 주소",
        "/api/packing",
        "/packing은 다른 주소",
    ])
    bullet_panel(slide, 8.85, 1.65, 3.75, 4.9, "문법", [
        "record 필드는 괄호 안에 선언",
        "entity 오타 주의",
        "List 타입은 ParameterizedTypeReference",
    ])

    # 8
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide)
    title(slide, "다음 학습 지점", "ch03-prompt는 프롬프트 전략 전시관입니다", "NEXT")
    bullet_panel(slide, 0.75, 1.75, 5.8, 4.75, "ch03-prompt", [
        "Prompt Template",
        "Zero-shot / Few-shot",
        "Role Assignment",
        "Multi Messages",
        "Step-back / Self-consistency",
    ])
    bullet_panel(slide, 6.8, 1.75, 5.75, 4.75, "학습 목표", [
        "같은 ChatClient라도 질문 전략에 따라 결과가 달라짐",
        "기능 구현보다 프롬프트 설계 감각을 키움",
        "Day2 구조화 응답과 결합하면 실전 기능이 됨",
    ])

    prs.save(OUT_MAIN)
    prs.save(OUT_STYLE)


if __name__ == "__main__":
    make_deck()
    print(OUT_MAIN)
    print(OUT_STYLE)
